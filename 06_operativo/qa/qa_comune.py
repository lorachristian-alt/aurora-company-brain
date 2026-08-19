# -*- coding: utf-8 -*-
"""qa_comune — fondamenta condivise della suite QA delle note.

Non e' uno dei quattro controlli: e' cio' che tutti e quattro usano.
Contiene l'estrattore di testo congelato (metodo_01 §5-bis, riusato IDENTICO),
il caricamento delle note col loro frontmatter, la normalizzazione con la
tabella alias (metodo_03 §7.1 clausola 1) e la gestione di perimetro e report.

Non modifica MAI una nota: la QA riporta, non corregge (metodo_03 §7).
"""
import io, json, os, re, sys, unicodedata
from datetime import date

# La console di Windows apre in cp1252 e va in errore sul primo carattere fuori tabella —
# e i rilievi di questa suite citano testo dei grezzi, che di caratteri fuori tabella e'
# pieno. Si forza UTF-8 in uscita, tollerando i caratteri non rappresentabili.
for _flusso in (sys.stdout, sys.stderr):
    try:
        _flusso.reconfigure(encoding="utf-8", errors="replace")
    except (AttributeError, ValueError):
        pass

# ---------------------------------------------------------------- percorsi

REPO = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))          # 06_operativo\
RADICE = os.path.dirname(REPO)                                              # radice repo
VAULT = r"C:\Users\buulo\Desktop\aurora-cervello"
SOURCES = os.path.join(VAULT, "sources")
MANIFEST = os.path.join(REPO, "manifest_corpus_v1.1.json")
ALIAS_MD = os.path.join(RADICE, "01_metodo", "alias_entita.md")
QA_DIR = os.path.join(REPO, "qa")

CARTELLE = ["self", "areas", "projects", "docs", "entities", "concepts",
            "data", "outputs", "code", "workspace", "sources"]

# workspace\ e sources\ sono esclusi dai conteggi di qualita' (metodo_03 §7.0)
ESCLUSE_QUALITA = {"workspace", "sources"}

# metodo_03 §2.4 — LA NOTA-STRUMENTO DEL PROGETTO, definita una volta sola.
#
# Emendamento E1 (gate del 16/08/2026) esteso da E20 (gate della matrice, 18/08/2026):
# questa e' l'UNICA definizione della classe, e tutte e tre le esenzioni si riferiscono a
# lei. Se un giorno se ne aggiunge una quarta, si aggiunge qui e non altrove: due
# definizioni della stessa classe divergono in un mese.
#
# Una nota-strumento documenta un ATTREZZO DEL PROGETTO — uno script della suite QA, un
# generatore di derivati — che non discende da nessun grezzo del corpus e non afferma
# nessun fatto di Aurora.
#
# ⚠️ La classe e' definita dal PREFISSO `script-` dentro `code\`, non dalla cartella: una
# nota che documenta un'automazione AZIENDALE (l'OCR dei DDT, l'integrazione EDI-ERP)
# parla di un fatto di Aurora, ha grezzi che la attestano, resta a schema pieno e resta
# soggetta a tutti i controlli, componente unica compresa. Se una nota di contenuto di
# `code\` e' staccata dal grafo, e' un difetto vero.
CARTELLA_STRUMENTI = "code"
PREFISSO_STRUMENTO = "script-"


def e_nota_strumento(nota):
    """Vero per le note che documentano un attrezzo del progetto (§2.4, E1 + E20).

    Le TRE esenzioni della classe, e nessun'altra:
      1. `fonti` e il blocco `## Fonti` sono facoltativi (§2.4, E1);
      2. restano fuori dallo strato di giudizio della provenance (§7.1 clausola 6);
      3. restano fuori dal controllo di COMPONENTE UNICA (§7.2, E20).

    Restano invece soggette, senza sconti, a: schema del frontmatter, wikilink rotti,
    nomi ambigui, e raggiungibilita' da `_index-code` (non sono orfane per esenzione).
    Si rivedono a occhio a ogni gate: e' l'unico controllo di merito che le riguarda."""
    return nota.cartella == CARTELLA_STRUMENTI and nota.slug.startswith(PREFISSO_STRUMENTO)

TYPE_AMMESSI = {"atomica", "hub", "entita", "conflitto", "concetto",
                "index", "sessione", "daily"}

AREE = {"qualita", "produzione", "manutenzione", "commerciale", "logistica",
        "amministrazione", "risorse-umane", "sicurezza-ambiente",
        "ricerca-sviluppo", "direzione"}

PREFISSI = ["self-", "area-", "fatto-", "progetto-", "entita-", "marchio-",
            "macchina-", "prodotto-", "lotto-", "doc-", "concetto-", "kpi-",
            "output-", "questione-", "script-", "bozza-", "sessione-",
            "diario-", "_index-"]

ESITO_ERRORE, ESITO_AVVISO = "ERRORE", "AVVISO"


# ------------------------------------------------- estrattore congelato

def text_of(p):
    """metodo_01 §5-bis. Riusata IDENTICA: se cambia qui, i confronti della
    provenance non sono piu' quelli con cui sono state verificate le 282 risposte.
    Nessun ramo .jpg: su una foto restituisce stringa vuota, ed e' il motivo
    per cui esiste il campo `verifica: visiva` (metodo_03 §2.3)."""
    e = p.rsplit(".", 1)[-1].lower()
    if e in ("txt", "csv", "log", "xml", "md", "p7m"):
        raw = open(p, "rb").read()
        for enc in ("utf-8", "cp1252"):
            try: return raw.decode(enc)
            except UnicodeDecodeError: pass
        return raw.decode("latin-1")
    if e == "pdf":
        from pypdf import PdfReader
        return "\n".join((pg.extract_text() or "") for pg in PdfReader(p).pages)
    if e == "docx":
        from docx import Document
        d = Document(p)
        return "\n".join([q.text for q in d.paragraphs] +
                         [c.text for t in d.tables for r in t.rows for c in r.cells])
    if e == "xlsx":
        from openpyxl import load_workbook
        wb = load_workbook(p, data_only=True)
        return "\n".join(" ".join(str(c) for c in row if c is not None)
                         for ws in wb.worksheets for row in ws.iter_rows(values_only=True))
    if e == "pptx":
        from pptx import Presentation
        out = []
        for sl in Presentation(p).slides:
            for sh in sl.shapes:
                if sh.has_text_frame: out.append(sh.text_frame.text)
                if sh.has_table:
                    out += [c.text for r in sh.table.rows for c in r.cells]
            if sl.has_notes_slide: out.append(sl.notes_slide.notes_text_frame.text)
        return "\n".join(out)
    if e == "eml":
        from email import policy
        from email.parser import BytesParser
        m = BytesParser(policy=policy.default).parse(open(p, "rb"))
        parti = ["%s: %s" % (h, m[h]) for h in
                 ("From", "To", "Cc", "Subject", "Date", "Message-ID") if m[h]]
        b = m.get_body(preferencelist=("plain", "html"))
        if b is not None: parti.append(b.get_content())
        parti += ["[allegato] %s" % a.get_filename()
                  for a in m.iter_attachments() if a.get_filename()]
        return "\n".join(parti)
    return ""


_cache_testo = {}

def testo_fonte(nome):
    """Testo estratto di un grezzo, letto una volta sola per esecuzione."""
    if nome not in _cache_testo:
        p = os.path.join(SOURCES, nome)
        try:
            _cache_testo[nome] = text_of(p) if os.path.isfile(p) else ""
        except Exception as ex:
            _cache_testo[nome] = ""
            print("  [avvertenza] estrazione fallita su %s: %s" % (nome, ex))
    return _cache_testo[nome]


# ------------------------------------------------------- normalizzazione

# ---- conteggio delle frasi di un `summary` ---------------------------------------
# ⚠️ FIX approvato dal coordinatore al gate del lotto 1B, con PERIMETRO CHIUSO.
# Il controllo «summary contiene piu' di una frase» contava i punti fermi, e in questo
# corpus sbagliava su ogni riassunto che contenesse una ragione sociale o un protocollo:
# «Frigotecnica Berica S.r.l.», «prot. VE-2026-3391», «Registro revisori n. 148223».
# Il fix ALLENTA un controllo, quindi non e' generico: vale solo per le abbreviazioni
# elencate qui sotto, e l'elenco si allunga **solo per un caso attestato nel corpus**,
# mai per prudenza. Il collaudo pianta un summary davvero multi-frase, con dentro queste
# stesse abbreviazioni, per dimostrare che il buco non si e' aperto.
ABBREVIAZIONI = frozenset((
    "s.r.l.", "s.p.a.", "s.n.c.", "s.a.s.",
    "prot.", "rev.", "n.", "nn.", "art.", "artt.", "pag.", "pagg.", "cod.", "cfr.",
    "sig.", "sig.ra", "dott.", "dott.ssa", "ing.", "rag.", "p.i.", "geom.",
    "es.", "ecc.", "min.", "max.", "ca.", "c.a.", "u.m.",
))

_RE_FINE_FRASE = re.compile(r"[.!?](?=\s|$)")


def conta_frasi(s):
    """Quante frasi contiene una stringa, saltando i punti di abbreviazione.

    Un punto che chiude una sigla dell'elenco ABBREVIAZIONI non chiude la frase.
    I punti interni di una sigla — la `S.` e la `r.` di `S.r.l.` — non sono mai
    seguiti da spazio, quindi non entrano nemmeno nel conteggio.
    """
    s = (s or "").rstrip()
    n = 0
    for m in _RE_FINE_FRASE.finditer(s):
        pezzi = s[:m.end()].split()
        if s[m.start()] == "." and pezzi and pezzi[-1].lower() in ABBREVIAZIONI:
            continue
        n += 1
    return n


def senza_accenti(s):
    return "".join(c for c in unicodedata.normalize("NFD", s)
                   if unicodedata.category(c) != "Mn")

def norm(s):
    """Forma di confronto: niente accenti, minuscolo, spazi collassati,
    apostrofi e virgolette uniformati. Serve perche' il corpus ha encoding
    misti e trascrizioni senza accenti."""
    s = senza_accenti(s).lower()
    # il quoting delle mail spezza le frasi: \u00ab...iniezione\n> azoto...\u00bb non e' una frase
    # diversa da \u00ab...iniezione azoto...\u00bb, ed e' la forma in cui meta' delle citazioni di
    # questo corpus vive, perche' le mail arrivano inoltrate e citate a piu' livelli.
    s = re.sub(r"\n[ \t]*>+[ \t]*", "\n", s)
    s = s.replace("\u2019", "'").replace("\u2018", "'")
    s = s.replace("\u00ab", '"').replace("\u00bb", '"')
    s = s.replace("\u201c", '"').replace("\u201d", '"')
    s = re.sub(r"[\s\u00a0]+", " ", s)
    return s.strip()


def carica_alias():
    """Legge le tabelle della CLASSE A di alias_entita.md.

    La tabella non si ricopia qui dentro: il padrone e' `01_metodo\\alias_entita.md`
    (metodo_03 §6.2), che cresce a ogni lotto. Questo parser la legge a ogni
    esecuzione, cosi' una riga aggiunta la' vale subito qui.

    Restituisce una lista di GRUPPI di forme equivalenti. Le classi B e C non si
    leggono: unire i loro soggetti distruggerebbe una trappola voluta.
    """
    if not os.path.isfile(ALIAS_MD):
        return []
    testo = open(ALIAS_MD, encoding="utf-8").read()
    # si prende solo cio' che sta fra "## Classe A" e "## Classe B"
    m = re.search(r"^## Classe A(.*?)^## Classe B", testo, re.S | re.M)
    if not m:
        return []
    gruppi = []
    for riga in m.group(1).splitlines():
        riga = riga.strip()
        if not riga.startswith("|") or riga.startswith("|---"):
            continue
        celle = [c.strip() for c in riga.strip("|").split("|")]
        if len(celle) < 2:
            continue
        prima, seconda = celle[0], celle[1]
        if prima.lower().startswith(("nota padrona", "nel grezzo")):
            continue                                   # intestazione
        forme = []
        for cella in (prima, seconda):
            for pezzo in cella.split("\u00b7"):        # separatore ' · '
                pezzo = pezzo.replace("**", "").replace("`", "")
                pezzo = pezzo.replace("\u00ab", "").replace("\u00bb", "")
                pezzo = re.sub(r"\(.*?\)", "", pezzo)  # via le glosse fra parentesi
                pezzo = pezzo.strip(" \"'")
                if pezzo and len(pezzo) > 1 and not pezzo.startswith("["):
                    forme.append(pezzo)
        if len(forme) >= 2:
            gruppi.append([norm(f) for f in forme])
    return gruppi


_ALIAS = None

def gruppi_alias():
    global _ALIAS
    if _ALIAS is None:
        _ALIAS = carica_alias()
    return _ALIAS


# sostituzioni note dello scanner (alias_entita.md, «Come si legge»)
def varianti_ocr(t):
    """Genera le forme che lo scanner puo' aver prodotto: 0<->O, l<->1, S<->5."""
    fuori = {t}
    coppie = [("0", "o"), ("l", "1"), ("s", "5")]
    for a, b in coppie:
        nuove = set()
        for f in fuori:
            nuove.add(f.replace(a, b)); nuove.add(f.replace(b, a))
        fuori |= nuove
    return fuori


def forme_equivalenti(tok):
    """Tutte le scritture sotto cui `tok` puo' comparire in un grezzo."""
    n = norm(tok)
    forme = {n}
    for g in gruppi_alias():
        if n in g:
            forme |= set(g)
    forme |= varianti_ocr(n)
    # un codice puo' comparire con o senza trattini e spazi
    compatto = re.sub(r"[-\s/_.]", "", n)
    if len(compatto) > 3:
        forme.add(compatto)
    return forme


def presente(tok, testo_norm, testo_compatto):
    """Il token compare nel testo, in una qualunque delle sue forme?"""
    for f in forme_equivalenti(tok):
        if not f:
            continue
        if f in testo_norm:
            return True
        fc = re.sub(r"[-\s/_.]", "", f)
        if len(fc) > 3 and fc in testo_compatto:
            return True
    return False


# --------------------------------------------------------- lettura note

class Nota(object):
    def __init__(self, percorso):
        self.percorso = percorso
        self.nome = os.path.basename(percorso)
        self.slug = self.nome[:-3]
        self.cartella = os.path.basename(os.path.dirname(percorso))
        self.grezzo = open(percorso, encoding="utf-8").read()
        self.fm, self.corpo, self.errore_fm, self.riga_fm = _spezza(self.grezzo)

    # --- comodita' -------------------------------------------------
    @property
    def type(self):   return (self.fm or {}).get("type")
    @property
    def area(self):   return (self.fm or {}).get("area")
    @property
    def stato(self):  return (self.fm or {}).get("stato")
    @property
    def fonti(self):
        f = (self.fm or {}).get("fonti")
        if f is None: return []
        return f if isinstance(f, list) else [f]

    @property
    def corpo_senza_fonti(self):
        """Il corpo escluso il blocco `## Fonti`: e' apparato, non contenuto."""
        return re.split(r"^##\s+Fonti\s*$", self.corpo, flags=re.M)[0]

    @property
    def blocco_fonti(self):
        p = re.split(r"^##\s+Fonti\s*$", self.corpo, flags=re.M)
        return p[1] if len(p) > 1 else ""

    def wikilink(self, solo_corpo=True):
        """I wikilink del CORPO. Quelli verso sources\\ non sono relazioni: sono fonti."""
        testo = self.corpo if solo_corpo else self.grezzo
        out = []
        for m in re.finditer(r"\[\[([^\]|]+)(?:\|[^\]]*)?\]\]", testo):
            out.append((m.group(1).strip(), testo[:m.start()].count("\n") + 1))
        return out

    def parole_corpo(self):
        t = self.corpo_senza_fonti
        t = re.sub(r"^#.*$", "", t, flags=re.M)          # via i titoli
        t = re.sub(r"[|`\[\]#*_>-]", " ", t)
        return len([w for w in t.split() if any(c.isalnum() for c in w)])


def _spezza(grezzo):
    """Separa frontmatter YAML e corpo. Restituisce (dict, corpo, errore, riga)."""
    import yaml
    if not grezzo.startswith("---"):
        return None, grezzo, "la nota non comincia con un frontmatter", 1
    fine = re.search(r"^---\s*$", grezzo[3:], re.M)
    if not fine:
        return None, grezzo, "frontmatter aperto e mai chiuso", 1
    testa = grezzo[3:3 + fine.start()]
    corpo = grezzo[3 + fine.end():]
    # un secondo frontmatter piu' avanti e' un errore di schema
    if re.search(r"^---\s*$", corpo, re.M) and corpo.lstrip().startswith("---"):
        return None, corpo, "la nota ha piu' di un frontmatter", 1
    try:
        fm = yaml.safe_load(testa)
    except Exception as ex:
        return None, corpo, "frontmatter YAML non valido: %s" % str(ex).split("\n")[0], 2
    if not isinstance(fm, dict):
        return None, corpo, "il frontmatter non e' una mappa di campi", 2
    return fm, corpo, None, 0


def tutte_le_note(vault=VAULT):
    note = []
    for c in CARTELLE:
        d = os.path.join(vault, c)
        if not os.path.isdir(d):
            continue
        for n in sorted(os.listdir(d)):
            if n.lower().endswith(".md"):
                note.append(Nota(os.path.join(d, n)))
    return note


# ------------------------------------------------------------ perimetro

def aggiungi_argomenti(ap):
    ap.add_argument("--perimetro", nargs="+", required=True,
                    metavar="MODO",
                    help="'vault' oppure 'lotto <file1> <file2> ...' "
                         "(oppure 'lotto @elenco.txt')")
    ap.add_argument("--report", default=None,
                    help="cartella del report; default 06_operativo\\qa\\<data>_<lotto>\\")
    ap.add_argument("--vault", default=VAULT)


def leggi_perimetro(args):
    """Restituisce (modo, insieme_file_del_lotto).

    In modalita' lotto i controlli di COPERTURA, AREE POPOLATE e COMPONENTE UNICA
    si valutano solo sui file e sulle note del lotto (metodo_03 §7). Gli altri
    controlli sono identici: un lotto non si chiude ammorbidendo la QA.
    """
    p = args.perimetro
    modo = p[0].lower()
    if modo == "vault":
        return "vault", None
    if modo != "lotto":
        print("perimetro sconosciuto: %s (attesi 'vault' o 'lotto')" % modo)
        sys.exit(1)
    voci = p[1:]
    if len(voci) == 1 and voci[0].startswith("@"):
        elenco = voci[0][1:]
        voci = [r.strip() for r in open(elenco, encoding="utf-8")
                if r.strip() and not r.strip().startswith("#")]
    if not voci:
        print("--perimetro lotto richiede l'elenco dei grezzi del lotto")
        sys.exit(1)
    return "lotto", set(voci)


def note_del_perimetro(note, modo, file_lotto):
    """Le note che il lotto ha prodotto: quelle che citano almeno un suo grezzo,
    piu' gli `_index` e gli hub d'area, che il lotto tocca sempre."""
    if modo == "vault":
        return note
    dentro = []
    for n in note:
        if n.type == "index" or any(f in file_lotto for f in n.fonti):
            dentro.append(n)
    return dentro


def cartella_report(args, modo, nome_lotto="lotto"):
    if args.report:
        d = args.report
    else:
        d = os.path.join(QA_DIR, "%s_%s" % (date.today().isoformat(),
                                            "vault" if modo == "vault" else nome_lotto))
    os.makedirs(d, exist_ok=True)
    return d


# --------------------------------------------------------------- report

class Report(object):
    """Raccoglie gli esiti e decide il codice di uscita.
    0 tutto verde · 1 almeno un ERRORE · 2 solo AVVISI (metodo_03 §7)."""

    def __init__(self, nome):
        self.nome = nome
        self.voci = []          # (esito, nota, riga, controllo, messaggio)

    def errore(self, nota, controllo, messaggio, riga=0):
        self.voci.append((ESITO_ERRORE, nota, riga, controllo, messaggio))

    def avviso(self, nota, controllo, messaggio, riga=0):
        self.voci.append((ESITO_AVVISO, nota, riga, controllo, messaggio))

    @property
    def errori(self): return [v for v in self.voci if v[0] == ESITO_ERRORE]
    @property
    def avvisi(self): return [v for v in self.voci if v[0] == ESITO_AVVISO]

    def codice(self):
        if self.errori: return 1
        if self.avvisi: return 2
        return 0

    def stampa(self):
        for esito, nota, riga, controllo, msg in self.voci:
            dove = "%s:%d" % (nota, riga) if riga else nota
            print("%-7s %-28s %-26s %s" % (esito, dove, controllo, msg))
        print("-- %s: %d errori, %d avvisi" % (self.nome, len(self.errori), len(self.avvisi)))

    def markdown(self):
        r = ["## %s\n" % self.nome,
             "- ERRORI: **%d**" % len(self.errori),
             "- AVVISI: **%d**\n" % len(self.avvisi)]
        for etichetta, voci in (("Errori", self.errori), ("Avvisi", self.avvisi)):
            if not voci:
                continue
            r.append("### %s\n" % etichetta)
            r.append("| Nota | Riga | Controllo | Rilievo |")
            r.append("|---|---|---|---|")
            for _, nota, riga, controllo, msg in voci:
                msg = msg.replace("|", "\\|")
                r.append("| `%s` | %s | %s | %s |" % (nota, riga or "", controllo, msg))
            r.append("")
        return "\n".join(r)


def scrivi_report(cartella, nomefile, testo):
    p = os.path.join(cartella, nomefile)
    with io.open(p, "w", encoding="utf-8") as f:
        f.write(testo)
    return p


def manifest_nomi():
    man = json.load(open(MANIFEST, encoding="utf-8"))
    return {e["nome"] for e in man["file"]}
