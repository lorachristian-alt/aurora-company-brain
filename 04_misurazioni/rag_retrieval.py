# -*- coding: utf-8 -*-
"""
rag_retrieval.py — Misura B, "configurazione B — RAG a embedding" (Config_test_13-08-26).

Costruisce SOLO il recupero: legge l'archivio `sources/`, lo spezza, lo indicizza con
embedding locali e per ogni domanda di `domande_solo.jsonl` estrae gli 8 passaggi piu'
vicini, con il nome del file di origine. Non risponde alle domande.

Parametri fissi dalla tabella di configurazione B (nessuno e' modificabile da riga di
comando: se cambiano, la misura non e' piu' confrontabile con la A):

    modello di embedding  BAAI/bge-m3, scaricato una volta e congelato in ./modelli/bge-m3
    estrazione testo      funzione `text_of` del §5-bis di docs_PMI_blueprint.md, identica
    chunking              1.200 caratteri, overlap 200, taglio su confine di riga
    vector store          Chroma persistente su disco, ./indice_B
    similarita'           coseno
    top_k                 8
    re-ranking            nessuno
    metadato per chunk    nome del file di origine (obbligatorio)

Riprendibilita': l'indicizzazione tiene traccia dei file gia' fatti in
`indice_B/stato_indice.json`; il recupero salta le domande gia' presenti in
`misuraB_contesti.jsonl` e scrive una riga alla volta con fsync. Un'interruzione in
qualunque punto costa al massimo il file (o la domanda) in corso.

Perimetro: legge `sources/` e `domande_solo.jsonl`, scrive solo dentro `misure_aurora/`.
Nessun file di risposte, soluzioni o valutazione viene aperto.
"""

import json
import os
import shutil
import sys
import time
from pathlib import Path

# ---------------------------------------------------------------- percorsi

BASE = Path(__file__).resolve().parent               # Desktop\misure_aurora
SOURCES = BASE.parent / "aurora-cervello" / "sources"
DOMANDE = BASE / "domande_solo.jsonl"
OUT = BASE / "misuraB_contesti.jsonl"
INDEX_DIR = BASE / "indice_B"
STATO = INDEX_DIR / "stato_indice.json"
MODEL_DIR = BASE / "modelli" / "bge-m3"

# ---------------------------------------------------------------- parametri

MODEL_NAME = "BAAI/bge-m3"
CHUNK = 1200
OVERLAP = 200
TOP_K = 8
COLLECTION = "misuraB_chunks"
MAX_SEQ = 1024          # un chunk da 1.200 caratteri non ci arriva mai: tetto di sicurezza
BATCH_EMB = 8           # batch di encoding su CPU
BATCH_ADD = 500         # righe per add() su Chroma

PARAMETRI = {           # se cambia uno di questi, l'indice va rifatto da zero
    "modello": MODEL_NAME,
    "chunk": CHUNK,
    "overlap": OVERLAP,
    "similarita": "cosine",
}


# ---------------------------------------------------------------- estrazione
# §5-bis di docs_PMI_blueprint.md, identica: e' il vincolo della configurazione B.

def text_of(p):
    e = p.rsplit(".", 1)[-1].lower()
    if e in ("txt", "csv", "log", "xml", "md", "p7m"):
        raw = open(p, "rb").read()
        for enc in ("utf-8", "cp1252"):
            try: return raw.decode(enc)
            except UnicodeDecodeError: pass
        return raw.decode("latin-1")
    if e == "pdf":
        from pypdf import PdfReader
        return "\n".join((pg.extract_text() or "") for pg in PdfReader(p).pages)
    if e == "docx":
        from docx import Document
        d = Document(p)
        return "\n".join([q.text for q in d.paragraphs] +
                         [c.text for t in d.tables for r in t.rows for c in r.cells])
    if e == "xlsx":
        from openpyxl import load_workbook
        wb = load_workbook(p, data_only=True)
        return "\n".join(" ".join(str(c) for c in row if c is not None)
                         for ws in wb.worksheets for row in ws.iter_rows(values_only=True))
    if e == "pptx":
        from pptx import Presentation
        out = []
        for sl in Presentation(p).slides:
            for sh in sl.shapes:
                if sh.has_text_frame: out.append(sh.text_frame.text)
                if sh.has_table:
                    out += [c.text for r in sh.table.rows for c in r.cells]
            if sl.has_notes_slide: out.append(sl.notes_slide.notes_text_frame.text)
        return "\n".join(out)
    if e == "eml":                      # senza questo ramo i controlli sulle mail sono ciechi
        from email import policy
        from email.parser import BytesParser
        m = BytesParser(policy=policy.default).parse(open(p, "rb"))
        parti = ["%s: %s" % (h, m[h]) for h in
                 ("From", "To", "Cc", "Subject", "Date", "Message-ID") if m[h]]
        b = m.get_body(preferencelist=("plain", "html"))
        if b is not None: parti.append(b.get_content())
        parti += ["[allegato] %s" % a.get_filename()
                  for a in m.iter_attachments() if a.get_filename()]
        return "\n".join(parti)
    return ""


# ---------------------------------------------------------------- chunking

def spezza(testo):
    """1.200 caratteri, overlap 200, taglio SOLO su confine di riga.

    L'archivio e' pieno di tabelle e registri: una riga non si spezza mai a meta', quindi
    un chunk puo' superare i 1.200 caratteri se una singola riga e' piu' lunga di cosi'.
    La coda di overlap e' fatta di righe intere per un massimo di 200 caratteri; se
    l'ultima riga del chunk e' da sola piu' lunga di 200, quella giuntura resta senza
    sovrapposizione (l'alternativa sarebbe tagliare dentro il record).
    """
    righe = testo.splitlines()
    pezzi, buf, lung = [], [], 0
    for r in righe:
        l = len(r) + 1
        if buf and lung + l > CHUNK:
            pezzi.append("\n".join(buf))
            coda, clen = [], 0
            for rr in reversed(buf):
                if clen + len(rr) + 1 > OVERLAP:
                    break
                coda.insert(0, rr)
                clen += len(rr) + 1
            if len(coda) >= len(buf):       # garantisce che si avanzi sempre
                coda = coda[1:]
            buf = list(coda)
            lung = sum(len(x) + 1 for x in buf)
        buf.append(r)
        lung += l
    if buf:
        pezzi.append("\n".join(buf))
    return [p for p in (x.strip() for x in pezzi) if p]


# ---------------------------------------------------------------- modello

def carica_modello():
    """Scarica bge-m3 la prima volta e lo congela su disco; poi legge sempre da li'."""
    from sentence_transformers import SentenceTransformer
    if MODEL_DIR.exists() and any(MODEL_DIR.iterdir()):
        print("modello: da disco (%s)" % MODEL_DIR)
        m = SentenceTransformer(str(MODEL_DIR), device="cpu")
    else:
        print("modello: primo scaricamento di %s (~2,3 GB), poi congelato in %s"
              % (MODEL_NAME, MODEL_DIR))
        m = SentenceTransformer(MODEL_NAME, device="cpu")
        MODEL_DIR.parent.mkdir(parents=True, exist_ok=True)
        try:
            m.save(str(MODEL_DIR))
        except Exception as exc:                       # noqa: BLE001
            print("  ATTENZIONE: non sono riuscito a salvarlo su disco (%s)" % exc)
    m.max_seq_length = MAX_SEQ
    return m


class Modello:
    """Carica bge-m3 solo quando serve davvero un vettore: se l'indice e' gia' completo e
    le domande sono gia' tutte fatte, il rerun non paga 2,2 GB per non fare niente."""

    def __init__(self):
        self._m = None

    def __call__(self):
        if self._m is None:
            self._m = carica_modello()
        return self._m


def vettori(model, testi, batch=BATCH_EMB):
    return model().encode(testi, batch_size=batch, normalize_embeddings=True,
                          show_progress_bar=False, convert_to_numpy=True).tolist()


# ---------------------------------------------------------------- indice

def apri_collezione():
    import chromadb
    client = chromadb.PersistentClient(path=str(INDEX_DIR))
    return client.get_or_create_collection(COLLECTION, metadata={"hnsw:space": "cosine"})


def leggi_stato():
    if STATO.exists():
        try:
            s = json.loads(STATO.read_text(encoding="utf-8"))
            if s.get("parametri") == PARAMETRI:
                return s
            print("parametri cambiati rispetto all'indice esistente: lo rifaccio da zero")
        except Exception:                              # noqa: BLE001
            print("stato dell'indice illeggibile: lo rifaccio da zero")
    return None


def scrivi_stato(s):
    """Su Windows chiunque tenga aperto il file (un editor, un `type`) fa fallire la
    sostituzione con WinError 5. Si riprova, e se proprio non si puo' si tira avanti: lo
    stato serve solo a riprendere, l'indice vero e' gia' su disco dentro Chroma."""
    INDEX_DIR.mkdir(parents=True, exist_ok=True)
    tmp = STATO.with_suffix(".tmp")
    tmp.write_text(json.dumps(s, ensure_ascii=False, indent=1), encoding="utf-8")
    for tentativo in range(6):
        try:
            tmp.replace(STATO)
            return
        except PermissionError:
            time.sleep(0.5 * (tentativo + 1))
    print("  ATTENZIONE: stato_indice.json bloccato, salto questo salvataggio")


def indicizza(model):
    """Indicizza `sources/` file per file. Riprende da dove si era interrotta."""
    stato = leggi_stato()
    if stato is None:
        shutil.rmtree(INDEX_DIR, ignore_errors=True)
        stato = {"parametri": PARAMETRI, "file": {}, "chunk": 0}

    col = apri_collezione()
    fatti = stato["file"]

    files = sorted(p for p in SOURCES.iterdir() if p.is_file())
    da_fare = []
    for p in files:
        st = p.stat()
        f = fatti.get(p.name)
        if f and f["size"] == st.st_size and abs(f["mtime"] - st.st_mtime) < 1e-6:
            continue
        da_fare.append(p)

    if not da_fare:
        print("indice gia' completo: %d chunk da %d file" % (col.count(), len(fatti)))
        return col

    print("indicizzazione: %d file da fare su %d (%d chunk gia' presenti)"
          % (len(da_fare), len(files), col.count()))
    t0 = time.time()
    for n, p in enumerate(da_fare, 1):
        try:
            testo = text_of(str(p))
        except Exception as exc:                       # noqa: BLE001
            print("  [%3d/%d] %-60s ILLEGGIBILE (%s)" % (n, len(da_fare), p.name,
                                                         type(exc).__name__))
            fatti[p.name] = {"size": p.stat().st_size, "mtime": p.stat().st_mtime,
                             "chunk": 0}          # es. il lock file di Word: non riprovarlo
            scrivi_stato(stato)
            continue
        pezzi = spezza(testo or "")
        if not pezzi:
            print("  [%3d/%d] %-60s nessun testo estraibile" % (n, len(da_fare), p.name))
            fatti[p.name] = {"size": p.stat().st_size, "mtime": p.stat().st_mtime,
                             "chunk": 0}
            scrivi_stato(stato)
            continue

        col.delete(where={"file": p.name})              # se il file era stato rifatto
        for i in range(0, len(pezzi), BATCH_ADD):
            lotto = pezzi[i:i + BATCH_ADD]
            col.add(ids=["%s#%04d" % (p.name, i + j) for j in range(len(lotto))],
                    embeddings=vettori(model, lotto),
                    documents=lotto,
                    metadatas=[{"file": p.name}] * len(lotto))

        fatti[p.name] = {"size": p.stat().st_size, "mtime": p.stat().st_mtime,
                         "chunk": len(pezzi)}
        stato["chunk"] = sum(f["chunk"] for f in fatti.values())
        scrivi_stato(stato)

        el = time.time() - t0
        print("  [%3d/%d] %-60s %4d chunk   %6.1fs   stima rimanente %5.1f min"
              % (n, len(da_fare), p.name, len(pezzi), el,
                 (el / n) * (len(da_fare) - n) / 60))

    print("indicizzazione finita: %d chunk totali in %.1f min"
          % (col.count(), (time.time() - t0) / 60))
    return col


# ---------------------------------------------------------------- recupero

def domande():
    out = []
    with open(DOMANDE, encoding="utf-8") as f:
        for riga in f:
            riga = riga.strip()
            if riga:
                d = json.loads(riga)
                out.append((d["id"], d["domanda"]))
    return out


def gia_fatte():
    ids = set()
    if not OUT.exists():
        return ids
    with open(OUT, encoding="utf-8") as f:
        for riga in f:
            riga = riga.strip()
            if not riga:
                continue
            try:
                ids.add(json.loads(riga)["id"])
            except Exception:                          # noqa: BLE001
                pass                                   # riga troncata da un'interruzione
    return ids


def recupera(model, col):
    tutte = domande()
    fatte = gia_fatte()
    restanti = [(i, d) for i, d in tutte if i not in fatte]
    if not restanti:
        print("recupero: tutte le %d domande sono gia' in %s" % (len(tutte), OUT.name))
        return
    print("recupero: %d domande da fare su %d" % (len(restanti), len(tutte)))

    emb = vettori(model, [d for _, d in restanti], batch=16)
    t0 = time.time()
    with open(OUT, "a", encoding="utf-8") as f:
        for n, ((qid, testo), v) in enumerate(zip(restanti, emb), 1):
            r = col.query(query_embeddings=[v], n_results=TOP_K,
                          include=["documents", "metadatas"])
            passaggi = [{"file": m["file"], "testo": t}
                        for m, t in zip(r["metadatas"][0], r["documents"][0])]
            f.write(json.dumps({"id": qid, "domanda": testo, "passaggi": passaggi},
                               ensure_ascii=False) + "\n")
            f.flush()
            os.fsync(f.fileno())
            if n % 25 == 0 or n == len(restanti):
                print("  %d/%d  (%.1fs)" % (n, len(restanti), time.time() - t0))
    print("scritto %s" % OUT)


# ---------------------------------------------------------------- main

def main():
    sys.stdout.reconfigure(line_buffering=True)   # il log serve anche se rediretto su file
    for p, che in ((SOURCES, "cartella sources"), (DOMANDE, "file delle domande")):
        if not p.exists():
            sys.exit("manca %s: %s" % (che, p))
    t0 = time.time()
    model = Modello()
    col = indicizza(model)
    recupera(model, col)
    n = sum(1 for _ in open(OUT, encoding="utf-8"))
    print("\n%s: %d righe. Totale %.1f min." % (OUT.name, n, (time.time() - t0) / 60))


if __name__ == "__main__":
    main()
