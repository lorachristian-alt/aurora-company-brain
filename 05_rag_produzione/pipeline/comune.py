# -*- coding: utf-8 -*-
"""
comune.py — le fondamenta della configurazione C: percorsi, config congelata,
estrazione del testo (con cache), chunking, tokenizzazione BM25, metadati.

Nessun parametro e' scritto qui dentro: TUTTI vengono da `config_c.json`, che e' la
configurazione congelata. Se un numero e' in questo file, e' un bug: si legge dal config
o non esiste. E' la regola che rende la Sessione 6 confrontabile con la baseline C.
"""

import hashlib
import json
import os
import re
import unicodedata
from pathlib import Path

# ------------------------------------------------------------------ percorsi

BASE = Path(__file__).resolve().parents[1]          # 05_rag_produzione\
RADICE = BASE.parent                                 # radice del repository
CONFIG = Path(os.environ.get("AURORA_CONFIG") or (BASE / "config_c.json"))

# `AURORA_LOCALE` sposta la cartella di lavoro pesante (indice, chunk, stato). Serve a due
# cose vere: collaudare il codice su un indice usa-e-getta senza toccare quello della
# misura, e dare alla Sessione 6 un indice NUOVO — quello della baseline non si riusa e
# non si aggiorna (metodo_02, addendum sul perimetro della misura «dopo»).
LOCALE = Path(os.environ.get("AURORA_LOCALE") or (BASE / "_locale_non_su_github"))
MODELLI = BASE / "_locale_non_su_github" / "modelli"
INDICE = LOCALE / "indice_qdrant"

# La cache dell'estrazione NON segue l'override: e' condivisa apposta fra le misure, ed
# e' proprio il suo riuso a garantire che la Sessione 6 legga lo stesso testo della
# baseline C. La chiave e' lo SHA-256 del file, quindi non puo' confondere due corpus.
CACHE_ESTRAZIONE = BASE / "_locale_non_su_github" / "cache_estrazione"


def carica_config(percorso=None):
    """La configurazione congelata. Si legge, non si modifica da codice."""
    # utf-8-sig e non utf-8: un editor di Windows puo' aggiungere il BOM al salvataggio,
    # e un config che non si apre per tre byte invisibili e' un pomeriggio buttato.
    return json.loads(Path(percorso or CONFIG).read_text(encoding="utf-8-sig"))


def sha256_file(percorso, blocco=1 << 20):
    h = hashlib.sha256()
    with open(percorso, "rb") as f:
        while True:
            b = f.read(blocco)
            if not b:
                break
            h.update(b)
    return h.hexdigest()


# ------------------------------------------------------------------ estrazione
# Il ramo testuale e' identico alla `text_of` del §5-bis di metodo_01, cioe' identico
# alla configurazione B: sui dieci formati non-immagine, B e C leggono gli stessi byte
# nello stesso modo, e la differenza fra le due misure resta l'architettura di recupero.
# L'undicesimo formato (.jpg) e' l'unica aggiunta di C: l'OCR, deciso il 17/08/2026.

FORMATI_TESTUALI = ("txt", "csv", "log", "xml", "md", "p7m")


def _testo_nativo(p):
    """La `text_of` congelata di metodo_01 §5-bis, copiata alla lettera."""
    e = p.rsplit(".", 1)[-1].lower()
    if e in FORMATI_TESTUALI:
        raw = open(p, "rb").read()
        for enc in ("utf-8", "cp1252"):
            try:
                return raw.decode(enc)
            except UnicodeDecodeError:
                pass
        return raw.decode("latin-1")
    if e == "pdf":
        from pypdf import PdfReader
        return "\n".join((pg.extract_text() or "") for pg in PdfReader(p).pages)
    if e == "docx":
        from docx import Document
        d = Document(p)
        return "\n".join([q.text for q in d.paragraphs] +
                         [c.text for t in d.tables for r in t.rows for c in r.cells])
    if e == "xlsx":
        from openpyxl import load_workbook
        wb = load_workbook(p, data_only=True)
        return "\n".join(" ".join(str(c) for c in row if c is not None)
                         for ws in wb.worksheets for row in ws.iter_rows(values_only=True))
    if e == "pptx":
        from pptx import Presentation
        out = []
        for sl in Presentation(p).slides:
            for sh in sl.shapes:
                if sh.has_text_frame:
                    out.append(sh.text_frame.text)
                if sh.has_table:
                    out += [c.text for r in sh.table.rows for c in r.cells]
            if sl.has_notes_slide:
                out.append(sl.notes_slide.notes_text_frame.text)
        return "\n".join(out)
    if e == "eml":                      # senza questo ramo i controlli sulle mail sono ciechi
        from email import policy
        from email.parser import BytesParser
        m = BytesParser(policy=policy.default).parse(open(p, "rb"))
        parti = ["%s: %s" % (h, m[h]) for h in
                 ("From", "To", "Cc", "Subject", "Date", "Message-ID") if m[h]]
        b = m.get_body(preferencelist=("plain", "html"))
        if b is not None:
            parti.append(b.get_content())
        parti += ["[allegato] %s" % a.get_filename()
                  for a in m.iter_attachments() if a.get_filename()]
        return "\n".join(parti)
    return ""


def _aggancia_tesseract(cfg):
    """Decide QUALE binario di tesseract si usa, in ordine di precedenza:
    variabile d'ambiente `AURORA_TESSERACT` → percorso del config, se esiste → PATH.
    Il percorso del config vince sul PATH perche' l'installer di Windows non aggiunge
    tesseract al PATH e affidarvisi renderebbe l'estrazione dipendente da come e' stato
    aperto il terminale."""
    import pytesseract
    scelto = os.environ.get("AURORA_TESSERACT")
    if not scelto:
        dal_config = cfg["estrazione"]["ocr"].get("eseguibile")
        if dal_config and Path(dal_config).exists():
            scelto = dal_config
    if scelto:
        pytesseract.pytesseract.tesseract_cmd = scelto
    return pytesseract


def versione_tesseract(cfg):
    """Versione dello strumento OCR: va a verbale, perche' un OCR diverso e' un corpus
    diverso. Restituisce None se tesseract non e' raggiungibile."""
    if not cfg["estrazione"]["ocr"]["attivo"]:
        return None
    try:
        return str(_aggancia_tesseract(cfg).get_tesseract_version())
    except Exception:                                   # noqa: BLE001
        return None


def _ocr(p, cfg):
    """OCR di un'immagine. Restituisce (testo, confidenza_media, n_parole).

    La soglia decide se il file e' una SCANSIONE DI DOCUMENTO (testo utile) o una
    FOTOGRAFIA (rumore): sopra soglia il testo entra nell'indice, sotto soglia il file
    entra con la sola scheda di metadati. Chunk-spazzatura nell'indice costano due volte:
    occupano il top-k e portano il generatore fuori strada.
    """
    from PIL import Image

    ocr = cfg["estrazione"]["ocr"]
    pytesseract = _aggancia_tesseract(cfg)

    with Image.open(p) as img:
        img = img.convert("L")          # scala di grigi: tesseract non usa il colore
        dati = pytesseract.image_to_data(img, lang=ocr["lingua"],
                                         config=ocr.get("parametri", ""),
                                         output_type=pytesseract.Output.DICT)

    parole, confidenze = [], []
    for testo, conf in zip(dati["text"], dati["conf"]):
        testo = (testo or "").strip()
        try:
            conf = float(conf)
        except (TypeError, ValueError):
            continue
        if testo and conf >= 0:
            parole.append(testo)
            confidenze.append(conf)

    testo = " ".join(parole)
    media = sum(confidenze) / len(confidenze) if confidenze else 0.0
    return testo, media, len(parole)


def estrai(percorso, cfg, usa_cache=True):
    """Testo di un file del corpus, con cache su disco.

    La cache ha per chiave lo SHA-256 del file grezzo: lo stesso byte-per-byte non si
    ri-estrae mai, nemmeno da una sessione diversa o con uno strumento aggiornato. E' la
    condizione perche' la misura «dopo» della Sessione 6 veda ESATTAMENTE il testo che ha
    visto la baseline C: un aggiornamento di tesseract fra agosto e ottobre cambierebbe il
    corpus senza che nessuno se ne accorga.

    Restituisce un dizionario: testo, origine (nativa|ocr|scheda), diagnostica dell'OCR.
    """
    p = Path(percorso)
    h = sha256_file(p)
    cache = CACHE_ESTRAZIONE / ("%s.json" % h)
    if usa_cache and cache.exists():
        return json.loads(cache.read_text(encoding="utf-8"))

    est = cfg["estrazione"]
    ext = p.suffix.lower().lstrip(".")
    voce = {"nome": p.name, "sha256": h, "formato": ext, "origine": "nativa",
            "ocr": None, "strumento": None}

    if ext in est["formati_immagine"]:
        if not est["ocr"]["attivo"]:
            voce.update(testo="", origine="scheda",
                        ocr={"motivo": "ocr disattivato nel config"})
        else:
            # Se l'OCR e' attivo ma lo strumento non c'e', si ALZA LA MANO. Trattarlo
            # come «nessun testo» finirebbe in cache, e da quel momento in poi il file
            # resterebbe cieco per sempre senza che nessuno se ne accorga.
            if versione_tesseract(cfg) is None:
                raise RuntimeError(
                    "OCR attivo nel config ma tesseract non e' raggiungibile. "
                    "Installalo (o indica il percorso in estrazione.ocr.eseguibile), "
                    "oppure metti estrazione.ocr.attivo = false e rifai l'indice.")
            testo, media, n = _ocr(str(p), cfg)
            soglie = est["ocr"]["soglia_testo_utile"]
            alfanum = sum(1 for c in testo if c.isalnum())
            utile = (alfanum >= soglie["min_caratteri_alfanumerici"]
                     and media >= soglie["min_confidenza_media"])
            voce.update(
                testo=testo if utile else "",
                origine="ocr" if utile else "scheda",
                ocr={"confidenza_media": round(media, 2), "parole": n,
                     "caratteri_alfanumerici": alfanum, "sopra_soglia": utile},
                strumento="tesseract %s" % (versione_tesseract(cfg) or "?"),
            )
    else:
        try:
            voce["testo"] = _testo_nativo(str(p))
        except Exception as exc:                        # noqa: BLE001
            # es. il lock file di Word: rumore realistico, non un guasto da riparare
            voce.update(testo="", origine="illeggibile",
                        ocr={"errore": type(exc).__name__})

    if usa_cache:
        CACHE_ESTRAZIONE.mkdir(parents=True, exist_ok=True)
        cache.write_text(json.dumps(voce, ensure_ascii=False), encoding="utf-8")
    return voce


# ------------------------------------------------------------------ chunking
# Identico alla configurazione B (metodo_02): 1.200 caratteri, overlap 200, taglio SOLO su
# confine di riga. E' l'unico modo perche' il confronto B/C isoli l'architettura e non il
# taglio dei pezzi.

def spezza(testo, cfg):
    """1.200 caratteri, overlap 200, taglio SOLO su confine di riga.

    L'archivio e' pieno di tabelle e registri: una riga non si spezza mai a meta', quindi
    un chunk puo' superare i 1.200 caratteri se una singola riga e' piu' lunga di cosi'.
    La coda di overlap e' fatta di righe intere per un massimo di 200 caratteri; se
    l'ultima riga del chunk e' da sola piu' lunga di 200, quella giuntura resta senza
    sovrapposizione (l'alternativa sarebbe tagliare dentro il record).
    """
    dim = cfg["chunking"]["caratteri"]
    ov = cfg["chunking"]["overlap"]
    righe = testo.splitlines()
    pezzi, buf, lung = [], [], 0
    for r in righe:
        l = len(r) + 1
        if buf and lung + l > dim:
            pezzi.append("\n".join(buf))
            coda, clen = [], 0
            for rr in reversed(buf):
                if clen + len(rr) + 1 > ov:
                    break
                coda.insert(0, rr)
                clen += len(rr) + 1
            if len(coda) >= len(buf):       # garantisce che si avanzi sempre
                coda = coda[1:]
            buf = list(coda)
            lung = sum(len(x) + 1 for x in buf)
        buf.append(r)
        lung += l
    if buf:
        pezzi.append("\n".join(buf))
    return [p for p in (x.strip() for x in pezzi) if p]


# ------------------------------------------------------------------ tokenizzazione BM25
# Il punto critico dichiarato in metodo_04 §2: il tokenizzatore sparso NON deve spezzare i
# codici. `L26130-L1-T2` deve restare cercabile intero E per le sue parti, altrimenti una
# domanda che cita `L26130` non trova il documento che scrive il codice per esteso.

_TOKEN = re.compile(r"[0-9a-zàèéìòùç]+"
                    r"(?:[-_./][0-9a-zàèéìòùç]+)*")


def normalizza(testo):
    """Minuscole e forma unicode stabile. Gli accenti NON si tolgono: in italiano
    distinguono parole («e» / «e'»), e toglierli sposterebbe l'IDF."""
    return unicodedata.normalize("NFC", testo).lower()


def tokenizza(testo, cfg=None):
    """Token per BM25: il composto intero PIU' le sue parti alfanumeriche.

    `MOD-QA-31` produce ['mod-qa-31', 'mod', 'qa', '31']. Costa qualche termine in piu'
    nell'indice e fa trovare il documento sia a chi scrive il codice per esteso sia a chi
    ne cita un pezzo. Nessuno stemming e nessuna stoplist: lo stemming italiano
    massacrerebbe i codici, e le parole vuote le spegne gia' l'IDF.
    """
    fuori = []
    for m in _TOKEN.finditer(normalizza(testo)):
        intero = m.group(0)
        fuori.append(intero)
        if any(c in intero for c in "-_./"):
            parti = [x for x in re.split(r"[-_./]", intero) if x]
            if len(parti) > 1:
                fuori.extend(parti)
    return fuori


# ------------------------------------------------------------------ metadati per chunk
# SOLO cio' che si ricava dal corpus. Il canone (`canone_aurora.md`) e la tabella alias
# NON entrano: iniettare nell'indice fatti che il corpus grezzo non contiene misurerebbe
# un archivio gia' in parte organizzato, cioe' esattamente cio' che la Sessione 6 deve
# misurare DOPO. Nel vault canonizzato questi metadati arriveranno dal frontmatter.

_CODICI = [
    re.compile(r"\bL\d{5}(?:-L\d(?:-T\d)?)?\b"),          # lotti: L26130, L26130-L1-T2
    re.compile(r"\b(?:NC|REC|RDA|DDT|ODP)-\d{4}-\d{2,4}\b"),
    re.compile(r"\b(?:MOD|IO|PRP|PRC)-[A-Z]{2,3}-\d{2,3}\b"),
    re.compile(r"\bMV\d{2}-\d{4}[A-Z]?\b"),
    re.compile(r"\bAF-[A-Z]{2}-\d{3,4}\b"),
    re.compile(r"\b[A-Z]{1,3}-\d{2,4}\b"),                # E-214, AL-217, CCP-2
]
_DATA = re.compile(r"\b(\d{2})[/.-](\d{2})[/.-](20\d{2})\b")


def metadati_chunk(testo, limite_codici=24):
    """Codici e date presenti nel pezzo. Servono come payload filtrabile in Qdrant e
    come appiglio di lettura nelle tracce di audit."""
    codici = []
    for rx in _CODICI:
        for m in rx.finditer(testo):
            v = m.group(0)
            if v not in codici:
                codici.append(v)
    date = []
    for g, m_, a in _DATA.findall(testo):
        iso = "%s-%s-%s" % (a, m_, g)
        if iso not in date:
            date.append(iso)
    return {"codici": sorted(codici)[:limite_codici], "date": sorted(date)[:limite_codici]}


# ------------------------------------------------------------------ utilita'

def fissa_thread(cfg):
    """Numero di thread fisso: non cambia i risultati, ma rende i tempi confrontabili
    fra una sessione e l'altra e evita che torch prenda tutti i core su una macchina
    che ne ha quattro veri."""
    n = cfg["esecuzione"]["thread_cpu"]
    os.environ.setdefault("OMP_NUM_THREADS", str(n))
    os.environ.setdefault("MKL_NUM_THREADS", str(n))
    try:
        import torch
        torch.set_num_threads(n)
    except ImportError:
        pass
